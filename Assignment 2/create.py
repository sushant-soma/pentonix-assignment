from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation
prs = Presentation()

# Define slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

# Define size variables for the image and text on the first slide
image_width = Inches(4)
image_height = Inches(3)  # Adjust the image height as needed
text_width = slide_width - image_width
text_height = slide_height

# Add the first slide
slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

# Add the image on the right side of the first slide
img_path1 = 'slide1.png'  # Replace with the path to your first image
left_img1 = slide_width - image_width
top_img1 = (slide_height - image_height) / 2  # Center the image vertically
slide1.shapes.add_picture(img_path1, left_img1, top_img1, image_width, image_height)

# Add the text on the left side of the first slide
left_text1 = 0
top_text1 = 0
text_box1 = slide1.shapes.add_textbox(left_text1, top_text1, text_width, text_height)
text_frame1 = text_box1.text_frame

p1 = text_frame1.add_paragraph()
p1.text = "Presentation Skills"
font1 = p1.runs[0].font
font1.color.rgb = RGBColor(0, 0, 255)
p1.alignment = PP_ALIGN.CENTER  # Center align the text

# Add the second slide with introduction text and image
slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

# Add introduction text to the second slide
left_intro = Inches(0.5)  # Adjust the left position as needed
top_intro = Inches(1.0)  # Adjust the top position as needed
width_intro = slide_width - Inches(1.0)
height_intro = Inches(2.0)  # Adjust the height as needed

text_box2 = slide2.shapes.add_textbox(left_intro, top_intro, width_intro, height_intro)
text_frame2 = text_box2.text_frame

p2 = text_frame2.add_paragraph()
p2.text = "Introduction"
font2 = p2.runs[0].font
font2.size = Pt(36)  # Adjust the font size as needed
font2.color.rgb = RGBColor(0, 0, 255)  # Set text color to blue
p2.alignment = PP_ALIGN.CENTER  # Center align the heading

# Add the image to the second slide
img_path2 = 'slide2.png'  # Replace with the path to your second image
image_width2 = slide_width - Inches(1.0)  # Adjust the width to fit within the slide
image_height2 = slide_height - top_intro - height_intro - Inches(1)  # Adjust the height as needed
left_img2 = Inches(0.5)  # Adjust the left position as needed
top_img2 = top_intro + height_intro + Inches(0.25)  # Position the image below the text with some spacing

slide2.shapes.add_picture(img_path2, left_img2, top_img2, image_width2, image_height2)


# Add the third slide with learning objectives text and image
slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

# Add learning objectives text to the third slide
left_lo = Inches(0.5)  # Adjust the left position as needed
top_lo = Inches(1.0)  # Adjust the top position as needed
width_lo = slide_width - Inches(1.0)
height_lo = Inches(2.0)  # Adjust the height as needed

text_box3 = slide3.shapes.add_textbox(left_lo, top_lo, width_lo, height_lo)
text_frame3 = text_box3.text_frame

p3 = text_frame3.add_paragraph()
p3.text = "Learning Objectives"
font3 = p3.runs[0].font
font3.size = Pt(36)  # Adjust the font size as needed
font3.color.rgb = RGBColor(0, 0, 255)  # Set text color to blue
p3.alignment = PP_ALIGN.CENTER  # Center align the heading

# Add the image to the third slide
img_path3 = 'slide3.png'  # Replace with the path to your third image
image_width3 = slide_width - Inches(1.0)  # Adjust the width to fit within the slide
image_height3 = slide_height - top_lo - height_lo - Inches(1)  # Adjust the height as needed
left_img3 = Inches(0.5)  # Adjust the left position as needed
top_img3 = top_lo + height_lo + Inches(0.25)  # Position the image below the text with some spacing

slide3.shapes.add_picture(img_path3, left_img3, top_img3, image_width3, image_height3)

# Add the fourth slide with image on the right and two texts on the left
slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

# Add the image on the right side of the fourth slide
img_path4 = 'slide4.png'  # Replace with the path to your fourth image
left_img4 = slide_width - image_width
top_img4 = (slide_height - image_height) / 2  # Center the image vertically
slide4.shapes.add_picture(img_path4, left_img4, top_img4, image_width, image_height)

# Add the heading text on the left side of the fourth slide
left_heading = Inches(0.5)  # Adjust the left position as needed
top_heading = Inches(1.0)  # Adjust the top position as needed
width_heading = text_width
height_heading = Inches(1.0)  # Adjust the height as needed

text_box_heading = slide4.shapes.add_textbox(left_heading, top_heading, width_heading, height_heading)
text_frame_heading = text_box_heading.text_frame

p_heading = text_frame_heading.add_paragraph()
p_heading.text = "The Four Types of Presentation"
font_heading = p_heading.runs[0].font
font_heading.color.rgb = RGBColor(0, 0, 255)  # Blue color
font_heading.size = Pt(36)  # Adjust the font size as needed
p_heading.alignment = PP_ALIGN.LEFT  # Left align the heading

# Add the subheading text below the heading on the left side
left_subheading = left_heading
top_subheading = top_heading + height_heading  # Position below the heading
width_subheading = text_width
height_subheading = Inches(1.0)  # Adjust the height as needed

text_box_subheading = slide4.shapes.add_textbox(left_subheading, top_subheading, width_subheading, height_subheading)
text_frame_subheading = text_box_subheading.text_frame

p_subheading = text_frame_subheading.add_paragraph()
p_subheading.text = "There are four kinds of presentations: informational, instructional, stimulating, and convincing."
font_subheading = p_subheading.runs[0].font
font_subheading.color.rgb = RGBColor(0, 0, 0)  # Black color
font_subheading.size = Pt(24)  # Adjust the font size as needed
p_subheading.alignment = PP_ALIGN.LEFT  # Left align the subheading


# Add the fifth slide with image on the right and two text boxes on the left
slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

# Add the image on the right side of the fifth slide
img_path5 = 'slide5.png'  # Replace with the path to your fifth image
left_img5 = slide_width - image_width
top_img5 = (slide_height - image_height) / 2  # Center the image vertically
slide5.shapes.add_picture(img_path5, left_img5, top_img5, image_width, image_height)

# Add the "About Us" text box on the left side
left_about_us = Inches(0.5)  # Adjust the left position as needed
top_about_us = Inches(1.0)  # Adjust the top position as needed
width_about_us = text_width
height_about_us = Inches(0.5)  # Adjust the height as needed

text_box_about_us = slide5.shapes.add_textbox(left_about_us, top_about_us, width_about_us, height_about_us)
text_frame_about_us = text_box_about_us.text_frame

p_about_us = text_frame_about_us.add_paragraph()
p_about_us.text = "About Us"
font_about_us = p_about_us.runs[0].font
font_about_us.color.rgb = RGBColor(255, 0, 0)  # Red color
font_about_us.size = Pt(18)  # Adjust the font size as needed

# Add the "Technical Skills" text box below "About Us"
left_tech_skills = left_about_us
top_tech_skills = top_about_us + height_about_us  # Position below "About Us"
width_tech_skills = text_width
height_tech_skills = Inches(0.5)  # Adjust the height as needed

text_box_tech_skills = slide5.shapes.add_textbox(left_tech_skills, top_tech_skills, width_tech_skills, height_tech_skills)
text_frame_tech_skills = text_box_tech_skills.text_frame

p_tech_skills = text_frame_tech_skills.add_paragraph()
p_tech_skills.text = "Technical Skills"
font_tech_skills = p_tech_skills.runs[0].font
font_tech_skills.color.rgb = RGBColor(0, 0, 255)  # Blue color
font_tech_skills.size = Pt(18)  # Adjust the font size as needed

# Add the sixth slide with image on the left and three text boxes on the right
slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

# Add the image on the left side of the sixth slide
img_path6 = 'slide6.png'  # Replace with the path to your sixth image
left_img6 = 0  # Image is on the left
top_img6 = (slide_height - image_height) / 2  # Center the image vertically
slide6.shapes.add_picture(img_path6, left_img6, top_img6, image_width, image_height)

# Add the "Assessments" text box on the right side
left_assessments = text_width  # Adjust the left position as needed
top_assessments = Inches(1.0)  # Adjust the top position as needed
width_assessments = slide_width - text_width
height_assessments = Inches(0.5)  # Adjust the height as needed

text_box_assessments = slide6.shapes.add_textbox(left_assessments, top_assessments, width_assessments, height_assessments)
text_frame_assessments = text_box_assessments.text_frame

p_assessments = text_frame_assessments.add_paragraph()
p_assessments.text = "Assessments"
font_assessments = p_assessments.runs[0].font
font_assessments.color.rgb = RGBColor(255, 0, 0)  # Red color
font_assessments.size = Pt(16)  # Adjust the font size as needed

# Add the "Check Your Understanding" text box below "Assessments"
left_cyu = left_assessments
top_cyu = top_assessments + height_assessments  # Position below "Assessments"
width_cyu = width_assessments
height_cyu = Inches(0.5)  # Adjust the height as needed

text_box_cyu = slide6.shapes.add_textbox(left_cyu, top_cyu, width_cyu, height_cyu)
text_frame_cyu = text_box_cyu.text_frame

p_cyu = text_frame_cyu.add_paragraph()
p_cyu.text = "Check Your Understanding"
font_cyu = p_cyu.runs[0].font
font_cyu.color.rgb = RGBColor(0, 0, 255)  # Blue color
font_cyu.size = Pt(20)  # Adjust the font size as needed

# Add the question text box below "Check Your Understanding"
left_question = left_cyu
top_question = top_cyu + height_cyu  # Position below "Check Your Understanding"
width_question = width_cyu
height_question = Inches(2.0)  # Adjust the height as needed

text_box_question = slide6.shapes.add_textbox(left_question, top_question, width_question, height_question)
text_frame_question = text_box_question.text_frame

p_question = text_frame_question.add_paragraph()
p_question.text = "Q1. Which of the following is the most important part of using presentation skills?\n" \
                 "A. Planning\n" \
                 "B. Preparation\n" \
                 "C. Visualizing success\n" \
                 "D. Exercise/drinking water beforehand"
font_question = p_question.runs[0].font
font_question.color.rgb = RGBColor(0, 0, 0)  # Black color
font_question.size = Pt(16)  # Adjust the font size as needed

# Save the presentation
prs.save('presentation.pptx')

# Add the seventh slide with "Key points" heading and image below
slide7 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

# Add the "Key points" heading at the upper center
left_key_points = Inches(0.5)  # Adjust the left position as needed
top_key_points = Inches(0.5)  # Adjust the top position as needed
width_key_points = slide_width - Inches(1.0)
height_key_points = Inches(1.0)  # Adjust the height as needed

text_box_key_points = slide7.shapes.add_textbox(left_key_points, top_key_points, width_key_points, height_key_points)
text_frame_key_points = text_box_key_points.text_frame

p_key_points = text_frame_key_points.add_paragraph()
p_key_points.text = "Key points"
font_key_points = p_key_points.runs[0].font
font_key_points.color.rgb = RGBColor(0, 0, 255)  # Blue color
font_key_points.size = Pt(36)  # Adjust the font size as needed
p_key_points.alignment = PP_ALIGN.CENTER  # Center align the heading

# Add the image below the heading
img_path7 = 'slide7.png'  # Replace with the path to your seventh image
image_width7 = slide_width - Inches(1.0)  # Adjust the width to fit within the slide
image_height7 = slide_height - top_key_points - height_key_points - Inches(1)  # Adjust the height as needed
left_img7 = Inches(0.5)  # Adjust the left position as needed
top_img7 = top_key_points + height_key_points + Inches(0.25)  # Position the image below the heading with some spacing

slide7.shapes.add_picture(img_path7, left_img7, top_img7, image_width7, image_height7)

# Add the eighth slide with image on the right and multiple text boxes and logos
slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

# Add the image on the right side of the eighth slide
img_path8 = 'slide8.png'  # Replace with the path to your eighth image
left_img8 = slide_width - image_width
top_img8 = (slide_height - image_height) / 2  # Center the image vertically
slide8.shapes.add_picture(img_path8, left_img8, top_img8, image_width, image_height)

# Add the "References" text box in red color and small size
left_references = Inches(0.5)  # Adjust the left position as needed
top_references = Inches(1.0)  # Adjust the top position as needed
width_references = slide_width - image_width
height_references = Inches(0.5)  # Adjust the height as needed

text_box_references = slide8.shapes.add_textbox(left_references, top_references, width_references, height_references)
text_frame_references = text_box_references.text_frame

p_references = text_frame_references.add_paragraph()
p_references.text = "References"
font_references = p_references.runs[0].font
font_references.color.rgb = RGBColor(255, 0, 0)  # Red color
font_references.size = Pt(14)  # Adjust the font size as needed

# Add the "Additional Resources" text box in blue color and bigger size
left_resources = left_references
top_resources = top_references + height_references  # Position below "References"
width_resources = width_references
height_resources = Inches(0.75)  # Adjust the height as needed

text_box_resources = slide8.shapes.add_textbox(left_resources, top_resources, width_resources, height_resources)
text_frame_resources = text_box_resources.text_frame

p_resources = text_frame_resources.add_paragraph()
p_resources.text = "Additional Resources"
font_resources = p_resources.runs[0].font
font_resources.color.rgb = RGBColor(0, 0, 255)  # Blue color
font_resources.size = Pt(18)  # Adjust the font size as needed

# Add the first logo "slide8_logo1.png" with "References" heading
img_path8_logo1 = 'slide8_logo1.png'  # Replace with the path to your first logo
image_width8_logo1 = Inches(1.0)  # Adjust the width of the logo
image_height8_logo1 = Inches(0.5)  # Adjust the height of the logo
left_img8_logo1 = left_references
top_img8_logo1 = top_resources + height_resources  # Position below "Additional Resources"
top_references_heading = top_img8_logo1 + image_height8_logo1 + Inches(0.1)  # Position the heading below the logo

slide8.shapes.add_picture(img_path8_logo1, left_img8_logo1, top_img8_logo1, image_width8_logo1, image_height8_logo1)

# Add the "References" heading in black color below the first logo
text_box_references_heading = slide8.shapes.add_textbox(left_references, top_references_heading, width_references, height_references)
text_frame_references_heading = text_box_references_heading.text_frame

p_references_heading = text_frame_references_heading.add_paragraph()
p_references_heading.text = "References"
font_references_heading = p_references_heading.runs[0].font
font_references_heading.color.rgb = RGBColor(0, 0, 0)  # Black color
font_references_heading.size = Pt(14)  # Adjust the font size as needed

# Add the second logo "slide8_logo2.png" with "Start Learning" heading
img_path8_logo2 = 'slide8_logo2.png'  # Replace with the path to your second logo
image_width8_logo2 = Inches(1.0)  # Adjust the width of the logo
image_height8_logo2 = Inches(0.5)  # Adjust the height of the logo
left_img8_logo2 = left_references
top_img8_logo2 = top_references_heading + height_references  # Position below the "References" heading
top_start_learning_heading = top_img8_logo2 + image_height8_logo2 + Inches(0.1)  # Position the heading below the logo

slide8.shapes.add_picture(img_path8_logo2, left_img8_logo2, top_img8_logo2, image_width8_logo2, image_height8_logo2)

# Add the "Start Learning" heading in black color below the second logo
text_box_start_learning_heading = slide8.shapes.add_textbox(left_references, top_start_learning_heading, width_references, height_references)
text_frame_start_learning_heading = text_box_start_learning_heading.text_frame

p_start_learning_heading = text_frame_start_learning_heading.add_paragraph()
p_start_learning_heading.text = "Start Learning"
font_start_learning_heading = p_start_learning_heading.runs[0].font
font_start_learning_heading.color.rgb = RGBColor(0, 0, 0)  # Black color
font_start_learning_heading.size = Pt(14)  # Adjust the font size as needed

# Add the ninth slide with text on the right and multiple text boxes and a button on the left
slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

# Add the "slide9.png" image on the right side
img_path9 = 'slide9.png'  # Replace with the path to your ninth image
left_img9 = slide_width - image_width
top_img9 = (slide_height - image_height) / 2  # Center the image vertically
slide9.shapes.add_picture(img_path9, left_img9, top_img9, image_width, image_height)

# Add the "feedback" text box in red color and small size
left_feedback = Inches(0.5)  # Adjust the left position as needed
top_feedback = Inches(1.0)  # Adjust the top position as needed
width_feedback = slide_width - image_width
height_feedback = Inches(0.5)  # Adjust the height as needed

text_box_feedback = slide9.shapes.add_textbox(left_feedback, top_feedback, width_feedback, height_feedback)
text_frame_feedback = text_box_feedback.text_frame

p_feedback = text_frame_feedback.add_paragraph()
p_feedback.text = "Feedback"
font_feedback = p_feedback.runs[0].font
font_feedback.color.rgb = RGBColor(255, 0, 0)  # Red color
font_feedback.size = Pt(14)  # Adjust the font size as needed

# Add the "Course Evaluation" text box in blue color and bigger size
left_evaluation = left_feedback
top_evaluation = top_feedback + height_feedback  # Position below "Feedback"
width_evaluation = width_feedback
height_evaluation = Inches(0.75)  # Adjust the height as needed

text_box_evaluation = slide9.shapes.add_textbox(left_evaluation, top_evaluation, width_evaluation, height_evaluation)
text_frame_evaluation = text_box_evaluation.text_frame

p_evaluation = text_frame_evaluation.add_paragraph()
p_evaluation.text = "Course Evaluation"
font_evaluation = p_evaluation.runs[0].font
font_evaluation.color.rgb = RGBColor(0, 0, 255)  # Blue color
font_evaluation.size = Pt(18)  # Adjust the font size as needed

# Add the "Kindly spare a few moments to fill out the course feedback form." text box in black color and small size
left_feedback_form = left_evaluation
top_feedback_form = top_evaluation + height_evaluation  # Position below "Course Evaluation"
width_feedback_form = width_evaluation
height_feedback_form = Inches(0.5)  # Adjust the height as needed

text_box_feedback_form = slide9.shapes.add_textbox(left_feedback_form, top_feedback_form, width_feedback_form, height_feedback_form)
text_frame_feedback_form = text_box_feedback_form.text_frame

p_feedback_form = text_frame_feedback_form.add_paragraph()
p_feedback_form.text = "Kindly spare a few moments to fill out the course feedback form."
font_feedback_form = p_feedback_form.runs[0].font
font_feedback_form.color.rgb = RGBColor(0, 0, 0)  # Black color
font_feedback_form.size = Pt(14)  # Adjust the font size as needed

# Add the "Submit" button on the left side
left_button = Inches(0.5)  # Adjust the left position as needed
top_button = top_feedback_form + height_feedback_form + Inches(0.5)  # Position below the text box
width_button = Inches(2.0)  # Adjust the width as needed
height_button = Inches(0.5)  # Adjust the height as needed

# Create a rounded rectangle shape to represent the button
button = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left_button, top_button, width_button, height_button)
button.fill.solid()
button.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue color for the button
button.text = "Submit"  # Button label
button.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White color for text
button.text_frame.paragraphs[0].font.size = Pt(14)  # Adjust the font size as needed
button.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align the text

# Save the presentation
prs.save('presentation.pptx')
